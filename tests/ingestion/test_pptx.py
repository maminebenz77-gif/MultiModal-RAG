import io
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from multimodal_rag.ingestion.pptx import parse_pptx
from multimodal_rag.ingestion.schema import ElementType
from multimodal_rag.providers.base import VisionProvider


class FakeVisionProvider(VisionProvider):
    def describe(self, image_bytes: bytes, prompt: str | None = None) -> str:
        return f"a description of {len(image_bytes)} bytes"


@pytest.fixture(autouse=True)
def _fake_vision(monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.setattr(
        "multimodal_rag.ingestion.vision.get_vision", lambda: FakeVisionProvider()
    )


def _tiny_png() -> bytes:
    from PIL import Image

    buffer = io.BytesIO()
    Image.new("RGB", (1, 1), color="blue").save(buffer, format="PNG")
    return buffer.getvalue()


_TINY_PNG = _tiny_png()


def _write_sample(tmp_path: Path) -> Path:
    presentation = Presentation()

    slide1 = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide1.shapes.title.text = "A Title"
    slide1.placeholders[1].text_frame.text = "A paragraph of body text."

    slide2 = presentation.slides.add_slide(presentation.slide_layouts[5])
    table_shape = slide2.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(3), Inches(1))
    table = table_shape.table
    table.cell(0, 0).text, table.cell(0, 1).text = "A", "B"
    table.cell(1, 0).text, table.cell(1, 1).text = "1", "2"

    slide3 = presentation.slides.add_slide(presentation.slide_layouts[5])
    chart_data = CategoryChartData()
    chart_data.categories = ["X", "Y"]
    chart_data.add_series("Series 1", (10, 20))
    slide3.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(3), Inches(2), chart_data
    )

    slide4 = presentation.slides.add_slide(presentation.slide_layouts[5])
    image_path = tmp_path / "pic.png"
    image_path.write_bytes(_TINY_PNG)
    slide4.shapes.add_picture(str(image_path), Inches(1), Inches(1))

    pptx_path = tmp_path / "deck.pptx"
    presentation.save(str(pptx_path))
    return pptx_path


def test_parses_title_paragraph_table_chart_and_image(tmp_path: Path) -> None:
    elements = parse_pptx(_write_sample(tmp_path))
    types = [el.type for el in elements]
    assert types == [
        ElementType.TITLE,
        ElementType.PARAGRAPH,
        ElementType.TABLE,
        ElementType.CHART,
        ElementType.IMAGE,
    ]


def test_slide_numbers_are_correct(tmp_path: Path) -> None:
    elements = parse_pptx(_write_sample(tmp_path))
    slides = [el.metadata.slide for el in elements]
    assert slides == [1, 1, 2, 3, 4]


def test_table_rendered_as_markdown(tmp_path: Path) -> None:
    elements = parse_pptx(_write_sample(tmp_path))
    table = next(el for el in elements if el.type == ElementType.TABLE)
    assert table.text == "| A | B |\n| --- | --- |\n| 1 | 2 |"


def test_native_chart_gets_exact_data_not_a_vision_guess(tmp_path: Path) -> None:
    elements = parse_pptx(_write_sample(tmp_path))
    chart = next(el for el in elements if el.type == ElementType.CHART)
    assert chart.description is not None
    assert "X: 10.0" in chart.description
    assert "Y: 20.0" in chart.description
    assert chart.description_status == "generated"
    assert chart.image_bytes is None


def test_picture_shape_goes_through_vision_describer(tmp_path: Path) -> None:
    elements = parse_pptx(_write_sample(tmp_path))
    image = next(el for el in elements if el.type == ElementType.IMAGE)
    assert image.image_bytes == _TINY_PNG
    assert image.description == f"a description of {len(_TINY_PNG)} bytes"
    assert image.description_status == "generated"
