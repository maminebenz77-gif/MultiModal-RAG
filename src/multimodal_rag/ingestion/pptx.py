"""PPTX parser: python-pptx -> common Element schema.

PPTX is the only format with a native, reliable chart signal
(`shape.has_chart`) — PDF/DOCX/Markdown can only ever guess "this image
is probably a chart" from a vision description. Native charts also expose
their underlying data (categories, series values) directly, which is
strictly more accurate than any vision caption could be (a vision model
can misread axis values; the chart's own data model cannot), so chart
elements skip get_vision() entirely and read their own data instead.
"""

from pathlib import Path
from typing import cast

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.shapes.base import BaseShape
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.picture import Picture

from .schema import Element, ElementMetadata, ElementType
from .tables import rows_to_markdown, summarize_table
from .vision import ImageDescriber

_TITLE_PLACEHOLDER_TYPES = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}


def parse_pptx(path: Path, summarize_tables: bool = False) -> list[Element]:
    presentation = Presentation(str(path))
    describer = ImageDescriber()

    elements: list[Element] = []
    position = 0

    for slide_number, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            element = _element_for_shape(
                shape, path, slide_number, position, describer, summarize_tables
            )
            if element is not None:
                elements.append(element)
                position += 1

    return elements


def _element_for_shape(
    shape: BaseShape,
    path: Path,
    slide_number: int,
    position: int,
    describer: ImageDescriber,
    summarize_tables: bool,
) -> Element | None:
    metadata = ElementMetadata(source_file=str(path), slide=slide_number, position=position)

    if shape.has_chart:
        return Element(
            type=ElementType.CHART,
            description=_describe_native_chart(shape),
            description_status="generated",
            metadata=metadata,
        )

    if shape.has_table:
        table = cast(GraphicFrame, shape).table
        rows = [[cell.text for cell in row.cells] for row in table.rows]
        markdown_table = rows_to_markdown(rows)
        summary = summarize_table(markdown_table) if summarize_tables else None
        return Element(
            type=ElementType.TABLE, text=markdown_table, table_summary=summary, metadata=metadata
        )

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        image_bytes = cast(Picture, shape).image.blob
        description, status = describer.describe(image_bytes)
        return Element(
            type=ElementType.IMAGE,
            image_bytes=image_bytes,
            description=description,
            description_status=status,
            metadata=metadata,
        )

    if shape.has_text_frame:
        # has_text_frame narrows this at runtime, but python-pptx's type
        # hierarchy doesn't expose a single common type to cast to here.
        text = shape.text_frame.text.strip()  # type: ignore[attr-defined]
        if not text:
            return None
        is_title = (
            shape.is_placeholder
            and shape.placeholder_format.type in _TITLE_PLACEHOLDER_TYPES
        )
        return Element(
            type=ElementType.TITLE if is_title else ElementType.PARAGRAPH,
            text=text,
            metadata=metadata,
        )

    return None


def _describe_native_chart(shape: BaseShape) -> str:
    chart = cast(GraphicFrame, shape).chart
    try:
        categories = [str(c) for c in chart.plots[0].categories]
    except Exception:
        categories = []

    lines = [f"Chart type: {chart.chart_type}"]
    try:
        for series in chart.series:
            values = list(series.values)
            if categories and len(categories) == len(values):
                pairs = ", ".join(f"{c}: {v}" for c, v in zip(categories, values, strict=True))
            else:
                pairs = ", ".join(str(v) for v in values)
            lines.append(f"Series '{series.name}': {pairs}")
    except Exception as exc:
        lines.append(f"[native chart data unavailable: {exc}]")

    return "\n".join(lines)
